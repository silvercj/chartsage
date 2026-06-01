"""Pure export-format builders for a ChartSage :class:`~schemas.Report`.

Each builder turns a validated ``Report`` (plus rendered chart PNGs or the
source CSV) into downloadable bytes/text for one format. They are **I/O-free**:
no Playwright, no network, no database — only in-memory work — so they are
fully unit-testable with fake image bytes.

The chart-image list (produced by ``pdf_export.render_chart_images``) has the
shape ``[{"chart_id": str, "png": bytes}, ...]`` in DOM order. Images are
matched to chart captions/titles by ``chart_id``, falling back to report order.

Public API:
    build_pptx(report, images)      -> bytes
    build_xlsx(report, source_csv)  -> bytes
    build_png_zip(images)           -> bytes
    build_markdown(report, images)  -> str
    build_html(report, images)      -> bytes
"""
from __future__ import annotations

import base64
import io
import re
import zipfile
from html import escape

from schemas import KeyMetric, Report

__all__ = [
    "build_pptx",
    "build_xlsx",
    "build_png_zip",
    "build_markdown",
    "build_html",
]


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
def _caption_for(report: Report, image: dict, index: int) -> str:
    """Caption for an image: match by chart_id, else fall back to report order."""
    chart_id = image.get("chart_id")
    for chart in report.charts:
        if chart.chart_id == chart_id:
            return chart.caption
    if 0 <= index < len(report.charts):
        return report.charts[index].caption
    return ""


def _title_for(report: Report, image: dict, index: int) -> str:
    """Chart title for an image: match by chart_id, else fall back to report order."""
    chart_id = image.get("chart_id")
    for chart in report.charts:
        if chart.chart_id == chart_id:
            return chart.spec.title
    if 0 <= index < len(report.charts):
        return report.charts[index].spec.title
    return chart_id or f"chart-{index + 1}"


def _first_sentence(text: str) -> str:
    """First sentence of ``text`` (up to the first period), trimmed."""
    text = (text or "").strip()
    if not text:
        return ""
    head = text.split(".", 1)[0].strip()
    return head or text


def _report_title(report: Report) -> str:
    """A display title: first sentence of the summary, else 'ChartSage Report'."""
    return _first_sentence(report.summary) or "ChartSage Report"


def _format_value(metric: KeyMetric) -> str:
    """Render a KPI value per its declared format."""
    value = metric.value
    if metric.format == "currency":
        return f"${value:,.2f}"
    if metric.format == "percent":
        return f"{value:,.1f}%"
    # number: drop a trailing ".0" so integers read cleanly
    if value == int(value):
        return f"{int(value):,}"
    return f"{value:,.2f}"


def _kpi_lines(report: Report) -> list[str]:
    """KPI label/value lines, one per key metric."""
    return [f"{m.label}: {_format_value(m)}" for m in report.key_metrics]


def _summary_paragraphs(report: Report) -> list[str]:
    """Summary split into non-empty paragraphs (blank-line separated)."""
    return [p.strip() for p in (report.summary or "").split("\n\n") if p.strip()]


_SLUG_RE = re.compile(r"[^A-Za-z0-9._-]+")


def _slug(text: str) -> str:
    """Filesystem-safe slug for a filename component."""
    text = (text or "").strip()
    text = text.replace(" ", "-")
    text = _SLUG_RE.sub("-", text).strip("-_.")
    return text or "chart"


# --------------------------------------------------------------------------- #
# PPTX
# --------------------------------------------------------------------------- #
def build_pptx(report: Report, images: list[dict]) -> bytes:
    """Build a PowerPoint deck: title, summary, KPI, then one slide per image.

    Slides:
      1. Title    — first sentence of the summary (or "ChartSage Report") + date.
      2. Summary  — the full summary text.
      3. KPIs     — one text line per ``report.key_metrics`` entry.
      4..N        — one slide per image: the chart picture + its caption.
    """
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt

    def _decoded(png: bytes):
        """Return (blob, width_px, height_px) for a PNG, decoding via Pillow.

        python-pptx must open the blob (Pillow) to detect its content-type and
        native size. Real chart screenshots decode fine and pass through
        unchanged; a bogus/corrupt blob is replaced with a small valid PNG
        placeholder so a single bad image never breaks the whole deck.
        """
        from PIL import Image as PILImage

        try:
            im = PILImage.open(io.BytesIO(png))
            im.load()
            return png, im.width, im.height
        except Exception:
            buf = io.BytesIO()
            PILImage.new("RGB", (1240, 760), (245, 245, 245)).save(buf, format="PNG")
            return buf.getvalue(), 1240, 760

    prs = Presentation()
    blank = prs.slide_layouts[6]  # fully blank layout — we add our own textboxes
    sw, sh = prs.slide_width, prs.slide_height

    def _add_textbox(slide, left, top, width, height):
        return slide.shapes.add_textbox(left, top, width, height).text_frame

    margin = Inches(0.5)
    content_w = sw - 2 * margin

    # --- Slide 1: Title ----------------------------------------------------- #
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, margin, Inches(2.2), content_w, Inches(2.0))
    tf.word_wrap = True
    tf.paragraphs[0].text = _report_title(report)
    tf.paragraphs[0].runs[0].font.size = Pt(40)
    tf.paragraphs[0].runs[0].font.bold = True
    sub = tf.add_paragraph()
    sub.text = report.generated_at
    if sub.runs:
        sub.runs[0].font.size = Pt(18)

    # --- Slide 2: Summary --------------------------------------------------- #
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, margin, margin, content_w, Inches(0.9))
    tf.paragraphs[0].text = "Summary"
    tf.paragraphs[0].runs[0].font.size = Pt(28)
    tf.paragraphs[0].runs[0].font.bold = True
    body = _add_textbox(slide, margin, Inches(1.4), content_w, sh - Inches(2.0))
    body.word_wrap = True
    paragraphs = _summary_paragraphs(report) or [report.summary or ""]
    body.paragraphs[0].text = paragraphs[0]
    for para in paragraphs[1:]:
        body.add_paragraph().text = para

    # --- Slide 3: KPIs ------------------------------------------------------ #
    slide = prs.slides.add_slide(blank)
    tf = _add_textbox(slide, margin, margin, content_w, Inches(0.9))
    tf.paragraphs[0].text = "Key Metrics"
    tf.paragraphs[0].runs[0].font.size = Pt(28)
    tf.paragraphs[0].runs[0].font.bold = True
    body = _add_textbox(slide, margin, Inches(1.4), content_w, sh - Inches(2.0))
    body.word_wrap = True
    kpi_lines = _kpi_lines(report)
    if kpi_lines:
        body.paragraphs[0].text = kpi_lines[0]
        for line in kpi_lines[1:]:
            body.add_paragraph().text = line
    else:
        body.paragraphs[0].text = "No key metrics."

    # --- Slides 4..N: one per chart image ----------------------------------- #
    for i, image in enumerate(images):
        slide = prs.slides.add_slide(blank)
        title = _title_for(report, image, i)
        head = _add_textbox(slide, margin, margin, content_w, Inches(0.8))
        head.word_wrap = True
        head.paragraphs[0].text = title
        head.paragraphs[0].runs[0].font.size = Pt(24)
        head.paragraphs[0].runs[0].font.bold = True

        blob, px_w, px_h = _decoded(image.get("png") or b"")
        # Place the picture in the band between the title and the caption,
        # scaled to fit while preserving aspect ratio. Passing explicit width
        # AND height keeps python-pptx from needing the image's native EMU size.
        pic_top = Inches(1.3)
        caption_h = Inches(1.0)
        avail_w = int(content_w)
        avail_h = int(sh - pic_top - caption_h - margin)
        # Assume 96 dpi → EMU; scale down to fit the band (never upscale).
        emu_per_px = 914400 / 96
        nat_w = px_w * emu_per_px
        nat_h = px_h * emu_per_px
        scale = min(avail_w / nat_w, avail_h / nat_h, 1.0)
        out_w = Emu(int(nat_w * scale))
        out_h = Emu(int(nat_h * scale))
        left = Emu(int((sw - out_w) / 2))
        slide.shapes.add_picture(io.BytesIO(blob), left, pic_top, width=out_w, height=out_h)

        cap = _caption_for(report, image, i)
        if cap:
            cap_tf = _add_textbox(slide, margin, sh - caption_h - margin, content_w, caption_h)
            cap_tf.word_wrap = True
            cap_tf.paragraphs[0].text = cap
            if cap_tf.paragraphs[0].runs:
                cap_tf.paragraphs[0].runs[0].font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# XLSX
# --------------------------------------------------------------------------- #
def build_xlsx(report: Report, source_csv: bytes) -> bytes:
    """Build a workbook with a 'Data' sheet (the source CSV) and a 'Summary' sheet.

    'Data'    — the source CSV parsed via pandas.
    'Summary' — KPI label/value rows, then a blank row, then the summary text.
    """
    import pandas as pd

    try:
        data_df = pd.read_csv(io.BytesIO(source_csv))
    except Exception:
        # Empty / unparseable CSV — still produce a (header-only) Data sheet.
        data_df = pd.DataFrame()

    # Summary sheet: KPI rows, a separator, then the narrative summary.
    summary_rows: list[tuple[str, object]] = [("Label", "Value")]
    for m in report.key_metrics:
        summary_rows.append((m.label, _format_value(m)))
    summary_rows.append(("", ""))
    summary_rows.append(("Summary", report.summary or ""))
    summary_df = pd.DataFrame(summary_rows)

    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        data_df.to_excel(writer, sheet_name="Data", index=False)
        summary_df.to_excel(writer, sheet_name="Summary", index=False, header=False)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# PNG zip
# --------------------------------------------------------------------------- #
def build_png_zip(images: list[dict], report: Report | None = None) -> bytes:
    """Zip every image as ``{index}-{safe-title-or-chart_id}.png`` (1-based)."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for i, image in enumerate(images):
            if report is not None:
                name_src = _title_for(report, image, i)
            else:
                name_src = image.get("chart_id") or f"chart-{i + 1}"
            filename = f"{i + 1}-{_slug(name_src)}.png"
            zf.writestr(filename, image.get("png") or b"")
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# Markdown
# --------------------------------------------------------------------------- #
def build_markdown(report: Report, images: list[dict]) -> str:
    """Build a Markdown document with the summary, KPIs, and per-chart captions+images."""
    lines: list[str] = []
    lines.append(f"# {_report_title(report)}")
    lines.append("")
    lines.append(f"_{report.generated_at}_")
    lines.append("")

    for para in _summary_paragraphs(report) or [report.summary or ""]:
        lines.append(para)
        lines.append("")

    kpi_lines = _kpi_lines(report)
    if kpi_lines:
        lines.append("**KPIs**")
        lines.append("")
        for line in kpi_lines:
            lines.append(f"- {line}")
        lines.append("")

    for i, image in enumerate(images):
        title = _title_for(report, image, i)
        caption = _caption_for(report, image, i)
        lines.append(f"## {title}")
        lines.append("")
        if caption:
            lines.append(caption)
            lines.append("")
        b64 = base64.b64encode(image.get("png") or b"").decode("ascii")
        lines.append(f"![{title}](data:image/png;base64,{b64})")
        lines.append("")

    if report.data_quality:
        lines.append("**Data quality**")
        lines.append("")
        for note in report.data_quality:
            lines.append(f"- {note}")
        lines.append("")

    return "\n".join(lines).rstrip("\n") + "\n"


# --------------------------------------------------------------------------- #
# HTML
# --------------------------------------------------------------------------- #
def build_html(report: Report, images: list[dict]) -> bytes:
    """Build a minimal standalone (utf-8) HTML doc; images as base64 ``<img>``."""
    title = _report_title(report)
    parts: list[str] = [
        "<!DOCTYPE html>",
        '<html lang="en">',
        "<head>",
        '<meta charset="utf-8">',
        '<meta name="viewport" content="width=device-width, initial-scale=1">',
        f"<title>{escape(title)}</title>",
        "<style>"
        "body{font-family:system-ui,-apple-system,Segoe UI,Roboto,sans-serif;"
        "max-width:900px;margin:2rem auto;padding:0 1rem;line-height:1.5;color:#111}"
        "img{max-width:100%;height:auto;display:block;margin:0.5rem 0}"
        "h1{margin-bottom:0.25rem}.date{color:#666;margin-top:0}"
        "figure{margin:1.5rem 0}figcaption{color:#444;font-size:0.95rem}"
        "</style>",
        "</head>",
        "<body>",
        f"<h1>{escape(title)}</h1>",
        f'<p class="date">{escape(report.generated_at)}</p>',
    ]

    for para in _summary_paragraphs(report) or [report.summary or ""]:
        parts.append(f"<p>{escape(para)}</p>")

    kpi_lines = _kpi_lines(report)
    if kpi_lines:
        parts.append("<h2>KPIs</h2>")
        parts.append("<ul>")
        for line in kpi_lines:
            parts.append(f"<li>{escape(line)}</li>")
        parts.append("</ul>")

    for i, image in enumerate(images):
        img_title = _title_for(report, image, i)
        caption = _caption_for(report, image, i)
        b64 = base64.b64encode(image.get("png") or b"").decode("ascii")
        parts.append("<figure>")
        parts.append(f"<h2>{escape(img_title)}</h2>")
        parts.append(f'<img alt="{escape(img_title)}" src="data:image/png;base64,{b64}">')
        if caption:
            parts.append(f"<figcaption>{escape(caption)}</figcaption>")
        parts.append("</figure>")

    if report.data_quality:
        parts.append("<h2>Data quality</h2>")
        parts.append("<ul>")
        for note in report.data_quality:
            parts.append(f"<li>{escape(note)}</li>")
        parts.append("</ul>")

    parts.append("</body>")
    parts.append("</html>")
    return ("\n".join(parts) + "\n").encode("utf-8")
